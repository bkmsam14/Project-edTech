"""
Document parser utility for extracting text from uploaded files.
Supports: .pdf, .docx, .pptx, .txt, .md
"""
import io
import logging

logger = logging.getLogger(__name__)


def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n\n".join(text_parts)


def parse_docx(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text_parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
    return "\n\n".join(text_parts)


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
        if slide_texts:
            text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def parse_file(file_bytes: bytes, filename: str) -> str:
    """Auto-detect file type and parse accordingly."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "pptx": parse_pptx,
    }

    if ext in parsers:
        return parsers[ext](file_bytes)
    elif ext in ("txt", "md", "tex"):
        return file_bytes.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: .{ext}")
