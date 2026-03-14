"""
File parser utilities for extracting text from various document formats.
Supports: PDF, DOCX, PPTX, HTML, and plain text.
"""

import os
from typing import Optional
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def parse_pdf(file_path: str) -> str:
    """Extract text content from a PDF file."""
    try:
        reader = PdfReader(file_path)
        text_content = []

        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text.strip():
                text_content.append(f"=== Page {page_num} ===\n{page_text}")

        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error parsing PDF: {e}"


def parse_docx(file_path: str) -> str:
    """Extract text content from a DOCX file."""
    try:
        doc = Document(file_path)
        text_content = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_content.append(row_text)

        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error parsing DOCX: {e}"


def parse_pptx(file_path: str) -> str:
    """Extract text content from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if len(slide_text) > 1:  # More than just the header
                text_content.append("\n".join(slide_text))

        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error parsing PPTX: {e}"


def parse_text_file(file_path: str) -> str:
    """Extract text from a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as e:
            return f"Error reading text file: {e}"
    except Exception as e:
        return f"Error reading text file: {e}"


def parse_file(file_path: str, file_type: Optional[str] = None) -> str:
    """
    Parse a file and extract its text content.
    Auto-detects file type if not provided.
    """
    if not os.path.exists(file_path):
        return f"File not found: {file_path}"

    # Auto-detect file type from extension if not provided
    if not file_type:
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            ".pdf": "pdf",
            ".docx": "docx",
            ".doc": "docx",
            ".pptx": "pptx",
            ".ppt": "pptx",
            ".txt": "text",
        }
        file_type = type_map.get(ext, "unknown")

    # Parse based on file type
    if file_type == "pdf":
        return parse_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return parse_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return parse_pptx(file_path)
    elif file_type in ["text", "txt"]:
        return parse_text_file(file_path)
    else:
        return f"Unsupported file type: {file_type}"


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    """
    Split text into overlapping chunks for better RAG performance.

    Args:
        text: The text to chunk
        chunk_size: Maximum size of each chunk in characters
        overlap: Number of characters to overlap between chunks

    Returns:
        List of text chunks
    """
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size

        # Try to break at a paragraph or sentence boundary
        if end < len(text):
            # Look for paragraph break
            para_break = text.rfind("\n\n", start, end)
            if para_break > start + chunk_size // 2:
                end = para_break + 2
            else:
                # Look for sentence break
                sentence_break = text.rfind(". ", start, end)
                if sentence_break > start + chunk_size // 2:
                    end = sentence_break + 2

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - overlap

    return chunks
